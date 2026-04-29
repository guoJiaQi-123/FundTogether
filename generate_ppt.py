#!/usr/bin/env python3
"""Generate graduation defense PPT for FundTogether crowdfunding platform."""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Installing python-pptx...")
    os.system(f"{sys.executable} -m pip install python-pptx pillow")
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE


class PPTGenerator:
    def __init__(self, style='academic_white'):
        self.prs = Presentation()
        self.style = style
        self.colors = self._get_colors(style)
        self.fonts = self._get_fonts()

    def _get_colors(self, style):
        schemes = {
            'academic_white': {
                'primary': RGBColor(0, 51, 102),
                'secondary': RGBColor(70, 130, 180),
                'accent': RGBColor(204, 0, 0),
                'text': RGBColor(51, 51, 51),
                'bg': RGBColor(255, 255, 255),
                'light_bg': RGBColor(240, 245, 250),
                'highlight': RGBColor(0, 102, 153),
                'border': RGBColor(200, 210, 220),
            },
            'business_blue': {
                'primary': RGBColor(30, 60, 114),
                'secondary': RGBColor(70, 130, 180),
                'accent': RGBColor(255, 193, 7),
                'text': RGBColor(51, 51, 51),
                'bg': RGBColor(255, 255, 255),
                'light_bg': RGBColor(235, 240, 248),
                'highlight': RGBColor(42, 82, 152),
                'border': RGBColor(180, 200, 220),
            }
        }
        return schemes.get(style, schemes['academic_white'])

    def _get_fonts(self):
        return {
            'title': 'Microsoft YaHei',
            'body': 'Microsoft YaHei',
            'code': 'Consolas',
        }

    def _set_slide_bg(self, slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_title_bar(self, slide, title, subtitle=None):
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(1.3)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.colors['primary']
        title_bar.line.fill.background()

        if subtitle:
            txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.7))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = self.fonts['title']

            sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(8.8), Inches(0.4))
            sub_tf = sub_box.text_frame
            sub_tf.word_wrap = True
            p2 = sub_tf.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(200, 220, 240)
            p2.font.name = self.fonts['body']
        else:
            txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(8.8), Inches(0.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = self.fonts['title']

    def _add_page_number(self, slide, num, total):
        txBox = slide.shapes.add_textbox(Inches(9.0), Inches(7.1), Inches(0.8), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{num}/{total}"
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors['secondary']
        p.alignment = PP_ALIGN.RIGHT

    def add_title_slide(self, title, subtitle='', info_lines=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.colors['bg'])

        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.15)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.colors['primary']
        top_bar.line.fill.background()

        bottom_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.35), Inches(10), Inches(0.15)
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = self.colors['primary']
        bottom_bar.line.fill.background()

        left_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.8), Inches(0.08), Inches(1.5)
        )
        left_accent.fill.solid()
        left_accent.fill.fore_color.rgb = self.colors['accent']
        left_accent.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(8), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.font.name = self.fonts['title']

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = self.colors['secondary']
            p2.font.name = self.fonts['body']
            p2.space_before = Pt(12)

        if info_lines:
            info_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(8), Inches(2.5))
            info_tf = info_box.text_frame
            info_tf.word_wrap = True
            for i, line in enumerate(info_lines):
                p = info_tf.paragraphs[0] if i == 0 else info_tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.font.color.rgb = self.colors['text']
                p.font.name = self.fonts['body']
                p.space_after = Pt(8)

        return slide

    def add_content_slide(self, title, bullets, subtitle=None, page_num=None, total_pages=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.colors['bg'])
        self._add_title_bar(slide, title, subtitle)

        left, top, width, height = Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(bullet, tuple):
                p.text = f"  {bullet[0]}"
                p.font.size = Pt(bullet[1] if len(bullet) > 1 else 16)
                p.font.color.rgb = bullet[2] if len(bullet) > 2 else self.colors['text']
                p.font.bold = bullet[3] if len(bullet) > 3 else False
            else:
                p.text = f"\u25b8  {bullet}"
                p.font.size = Pt(16)
                p.font.color.rgb = self.colors['text']
            p.font.name = self.fonts['body']
            p.space_after = Pt(10)

        if page_num and total_pages:
            self._add_page_number(slide, page_num, total_pages)

        return slide

    def add_two_column_slide(self, title, left_title, left_content, right_title, right_content,
                              subtitle=None, page_num=None, total_pages=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.colors['bg'])
        self._add_title_bar(slide, title, subtitle)

        left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.4))
        lhf = left_header.text_frame
        lhp = lhf.paragraphs[0]
        lhp.text = left_title
        lhp.font.size = Pt(16)
        lhp.font.bold = True
        lhp.font.color.rgb = self.colors['highlight']
        lhp.font.name = self.fonts['body']

        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.3), Inches(4.5))
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        for i, item in enumerate(left_content):
            p = left_tf.paragraphs[0] if i == 0 else left_tf.add_paragraph()
            p.text = f"\u25b8  {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['text']
            p.font.name = self.fonts['body']
            p.space_after = Pt(8)

        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.95), Inches(1.5), Inches(0.03), Inches(5.0)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.colors['border']
        divider.line.fill.background()

        right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(0.4))
        rhf = right_header.text_frame
        rhp = rhf.paragraphs[0]
        rhp.text = right_title
        rhp.font.size = Pt(16)
        rhp.font.bold = True
        rhp.font.color.rgb = self.colors['highlight']
        rhp.font.name = self.fonts['body']

        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(4.3), Inches(4.5))
        right_tf = right_box.text_frame
        right_tf.word_wrap = True
        for i, item in enumerate(right_content):
            p = right_tf.paragraphs[0] if i == 0 else right_tf.add_paragraph()
            p.text = f"\u25b8  {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['text']
            p.font.name = self.fonts['body']
            p.space_after = Pt(8)

        if page_num and total_pages:
            self._add_page_number(slide, page_num, total_pages)

        return slide

    def add_table_slide(self, title, headers, rows, subtitle=None, page_num=None, total_pages=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.colors['bg'])
        self._add_title_bar(slide, title, subtitle)

        rows_count = len(rows) + 1
        cols_count = len(headers)
        left, top = Inches(0.5), Inches(1.8)
        width, height = Inches(9), Inches(4.5)
        table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
        table = table_shape.table

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors['primary']
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.name = self.fonts['body']
                paragraph.alignment = PP_ALIGN.CENTER

        for row_idx, row in enumerate(rows, 1):
            for col_idx, cell_text in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.colors['light_bg']
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = self.colors['text']
                    paragraph.font.name = self.fonts['body']
                    paragraph.alignment = PP_ALIGN.CENTER

        if page_num and total_pages:
            self._add_page_number(slide, page_num, total_pages)

        return slide

    def add_architecture_slide(self, title, layers, subtitle=None, page_num=None, total_pages=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.colors['bg'])
        self._add_title_bar(slide, title, subtitle)

        layer_colors = [
            RGBColor(0, 102, 153),
            RGBColor(42, 82, 152),
            RGBColor(70, 130, 180),
            RGBColor(100, 160, 200),
            RGBColor(150, 190, 220),
        ]

        start_y = 1.6
        layer_height = 1.0
        gap = 0.15

        for i, (layer_name, layer_desc) in enumerate(layers):
            y = start_y + i * (layer_height + gap)
            color_idx = min(i, len(layer_colors) - 1)

            rect = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(y), Inches(8.4), Inches(layer_height)
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = layer_colors[color_idx]
            rect.line.fill.background()

            name_box = slide.shapes.add_textbox(Inches(1.0), Inches(y + 0.1), Inches(2.5), Inches(0.4))
            name_tf = name_box.text_frame
            name_p = name_tf.paragraphs[0]
            name_p.text = layer_name
            name_p.font.size = Pt(16)
            name_p.font.bold = True
            name_p.font.color.rgb = RGBColor(255, 255, 255)
            name_p.font.name = self.fonts['body']

            desc_box = slide.shapes.add_textbox(Inches(3.5), Inches(y + 0.1), Inches(5.5), Inches(0.8))
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = layer_desc
            desc_p.font.size = Pt(13)
            desc_p.font.color.rgb = RGBColor(240, 245, 250)
            desc_p.font.name = self.fonts['body']

            if i < len(layers) - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.DOWN_ARROW,
                    Inches(4.8), Inches(y + layer_height), Inches(0.4), Inches(gap)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.colors['border']
                arrow.line.fill.background()

        if page_num and total_pages:
            self._add_page_number(slide, page_num, total_pages)

        return slide

    def add_highlight_slide(self, title, highlights, subtitle=None, page_num=None, total_pages=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.colors['bg'])
        self._add_title_bar(slide, title, subtitle)

        col_count = min(len(highlights), 3)
        card_width = 2.6
        gap = 0.3
        total_width = col_count * card_width + (col_count - 1) * gap
        start_x = (10 - total_width) / 2

        for i, (hl_title, hl_desc) in enumerate(highlights):
            col = i % col_count
            row = i // col_count
            x = start_x + col * (card_width + gap)
            y = 1.8 + row * 2.8

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(card_width), Inches(2.4)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.colors['light_bg']
            card.line.color.rgb = self.colors['border']
            card.line.width = Pt(1)

            num_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(0.5), Inches(0.4))
            num_tf = num_box.text_frame
            num_p = num_tf.paragraphs[0]
            num_p.text = f"0{i+1}"
            num_p.font.size = Pt(22)
            num_p.font.bold = True
            num_p.font.color.rgb = self.colors['accent']
            num_p.font.name = self.fonts['title']

            title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.6), Inches(card_width - 0.4), Inches(0.4))
            title_tf = title_box.text_frame
            title_tf.word_wrap = True
            title_p = title_tf.paragraphs[0]
            title_p.text = hl_title
            title_p.font.size = Pt(14)
            title_p.font.bold = True
            title_p.font.color.rgb = self.colors['primary']
            title_p.font.name = self.fonts['body']

            desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.1), Inches(card_width - 0.4), Inches(1.2))
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = hl_desc
            desc_p.font.size = Pt(11)
            desc_p.font.color.rgb = self.colors['text']
            desc_p.font.name = self.fonts['body']

        if page_num and total_pages:
            self._add_page_number(slide, page_num, total_pages)

        return slide

    def add_summary_slide(self, title, points, conclusion, page_num=None, total_pages=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.colors['bg'])
        self._add_title_bar(slide, title)

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, point in enumerate(points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u2714  {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['text']
            p.font.name = self.fonts['body']
            p.space_after = Pt(12)

        conclusion_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(5.4), Inches(8.4), Inches(1.2)
        )
        conclusion_box.fill.solid()
        conclusion_box.fill.fore_color.rgb = self.colors['primary']
        conclusion_box.line.fill.background()

        tf2 = conclusion_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = conclusion
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.font.name = self.fonts['body']
        p2.alignment = PP_ALIGN.CENTER

        if page_num and total_pages:
            self._add_page_number(slide, page_num, total_pages)

        return slide

    def add_thank_you_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.colors['primary'])

        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.15)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.colors['accent']
        top_bar.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "\u611f\u8c22\u5404\u4f4d\u8001\u5e08\uff01"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = self.fonts['title']
        p.alignment = PP_ALIGN.CENTER

        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        p2 = sub_tf.paragraphs[0]
        p2.text = "\u6b22\u8fce\u6307\u5bfc\u4e0e\u6279\u8bc4"
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(200, 220, 240)
        p2.font.name = self.fonts['body']
        p2.alignment = PP_ALIGN.CENTER

        return slide

    def save(self, output_path):
        self.prs.save(output_path)
        return output_path


def build_presentation():
    gen = PPTGenerator(style='academic_white')
    total_pages = 15

    # ===== Slide 1: Title =====
    gen.add_title_slide(
        title="\u5728\u7ebf\u4f17\u7b79\u7f51\u7ad9\u7684\u8bbe\u8ba1\u4e0e\u5b9e\u73b0",
        subtitle="FundTogether \u4f17\u7b79\u5e73\u53f0",
        info_lines=[
            "\u5b66\u3000\u3000\u9662\uff1a\u8f6f\u4ef6\u5b66\u9662",
            "\u4e13\u3000\u3000\u4e1a\uff1a\u8f6f\u4ef6\u5de5\u7a0b",
            "\u59d3\u3000\u3000\u540d\uff1a\u90ed\u5bb6\u65d7",
            "\u5b66\u3000\u3000\u53f7\uff1a2022005946",
            "\u6307\u5bfc\u6559\u5e08\uff1a\u674e\u91d1\u5c9a \u526f\u6559\u6388",
        ]
    )

    # ===== Slide 2: Research Background =====
    gen.add_content_slide(
        title="\u7814\u7a76\u80cc\u666f\u4e0e\u610f\u4e49",
        subtitle="\u4e3a\u4ec0\u4e48\u9009\u62e9\u5728\u7ebf\u4f17\u7b79\u4f5c\u4e3a\u7814\u7a76\u8bfe\u9898\uff1f",
        bullets=[
            "\u4f20\u7edf\u878d\u8d44\u65b9\u5f0f\u95e8\u69db\u9ad8\u3001\u5468\u671f\u957f\u3001\u6d41\u7a0b\u7e41\u7410\uff0c\u96be\u4ee5\u6ee1\u8db3\u4e2d\u5c0f\u578b\u521b\u65b0\u4e3b\u4f53\u7684\u7b79\u8d44\u9700\u6c42",
            "\u4f17\u7b79\u4f5c\u4e3a\u65b0\u578b\u878d\u8d44\u6a21\u5f0f\uff0c\u5177\u5907\u4f4e\u95e8\u69db\u3001\u5e7f\u8986\u76d6\u3001\u793e\u7fa4\u5316\u7684\u663e\u8457\u4f18\u52bf",
            "\u56fd\u5185\u4f17\u7b79\u5e73\u53f0\u666e\u904d\u5b58\u5728\u5ba1\u6838\u8584\u5f31\u3001\u8d44\u91d1\u6d41\u8f6c\u4e0d\u900f\u660e\u3001\u652f\u4ed8\u5408\u89c4\u98ce\u9669\u7b49\u95ee\u9898",
            "\u672c\u8bfe\u9898\u65e8\u5728\u8bbe\u8ba1\u5e76\u5b9e\u73b0\u4e00\u4e2a\u529f\u80fd\u5b8c\u6574\u3001\u8d44\u91d1\u95ed\u73af\u900f\u660e\u3001\u5b89\u5168\u53ef\u9760\u7684\u5728\u7ebf\u4f17\u7b79\u5e73\u53f0",
            "\u91cd\u70b9\u653b\u5173\uff1a\u8d44\u91d1\u6258\u7ba1\u6a21\u578b\u3001\u72b6\u6001\u673a\u6cbb\u7406\u3001\u5b9e\u65f6\u6d88\u606f\u63a8\u9001\u3001\u591a\u7aef\u9002\u914d",
        ],
        page_num=2, total_pages=total_pages
    )

    # ===== Slide 3: Technology Stack =====
    gen.add_table_slide(
        title="\u6280\u672f\u9009\u578b\u603b\u89c8",
        subtitle="\u57fa\u4e8e\u4f01\u4e1a\u7ea7\u4e3b\u6d41\u6280\u672f\u6808\u7684\u524d\u540e\u7aef\u5206\u79bb\u67b6\u6784",
        headers=["\u6280\u672f\u9886\u57df", "\u6280\u672f\u9009\u578b", "\u6838\u5fc3\u804c\u8d23"],
        rows=[
            ["\u524d\u7aef\u6846\u67b6", "Vue 3 + Vite", "\u54cd\u5e94\u5f0fSPA\u3001Composition API"],
            ["UI\u7ec4\u4ef6", "Element Plus + ECharts", "\u4f01\u4e1a\u7ea7\u7ec4\u4ef6 + \u6570\u636e\u53ef\u89c6\u5316"],
            ["\u540e\u7aef\u6846\u67b6", "Spring Boot 3.x", "RESTful\u670d\u52a1\u3001\u81ea\u52a8\u88c5\u914d"],
            ["ORM\u5c42", "MyBatis-Plus", "\u589e\u5f3aCRUD\u3001\u5206\u9875\u3001\u903b\u8f91\u5220\u9664"],
            ["\u6743\u9650\u5b89\u5168", "Spring Security + JWT", "\u65e0\u72b6\u6001RBAC\u3001\u4f4d\u63a9\u7801\u89d2\u8272\u6a21\u578b"],
            ["\u6570\u636e\u5e93", "MySQL 8.0 + Redis 7.x", "ACID\u4e8b\u52a1 + \u7f13\u5b58/\u5206\u5e03\u5f0f\u9501"],
            ["\u5b9e\u65f6\u901a\u4fe1", "WebSocket", "\u9879\u76ee\u52a8\u6001\u3001\u516c\u544a\u5b9e\u65f6\u63a8\u9001"],
            ["\u90e8\u7f72", "Docker Compose + Nginx", "\u4e00\u952e\u7f16\u6392\u3001\u53cd\u5411\u4ee3\u7406"],
        ],
        page_num=3, total_pages=total_pages
    )

    # ===== Slide 4: System Architecture =====
    gen.add_architecture_slide(
        title="\u7cfb\u7edf\u603b\u4f53\u67b6\u6784",
        subtitle="\u524d\u540e\u7aef\u5206\u79bb + \u591a\u5c42\u67b6\u6784\u8bbe\u8ba1",
        layers=[
            ("\u524d\u7aef\u5c42 (Presentation)", "Vue 3 + Element Plus + ECharts | \u54cd\u5e94\u5f0fSPA | WebSocket Client"),
            ("\u7f51\u5173\u5c42 (Gateway)", "Nginx \u53cd\u5411\u4ee3\u7406 | \u9759\u6001\u8d44\u6e90\u6258\u7ba1 | \u8bf7\u6c42\u8f6c\u53d1\u4e0e\u9650\u6d41"),
            ("\u670d\u52a1\u5c42 (Business)", "Spring Boot 3.x | Controller-Service-Mapper\u4e09\u5c42 | JWT\u9274\u6743 | AOP\u9650\u6d41"),
            ("\u6570\u636e\u5c42 (Data)", "MySQL 8.0 (ACID\u4e8b\u52a1) | Redis 7.x (\u7f13\u5b58/\u5206\u5e03\u5f0f\u9501) | MyBatis-Plus ORM"),
            ("\u57fa\u7840\u8bbe\u65bd (Infra)", "Docker Compose\u7f16\u6392 | \u6587\u4ef6\u5b58\u50a8 | WebSocket\u670d\u52a1\u7aef"),
        ],
        page_num=4, total_pages=total_pages
    )

    # ===== Slide 5: Database Design =====
    gen.add_two_column_slide(
        title="\u6570\u636e\u5e93\u8bbe\u8ba1",
        subtitle="\u6838\u5fc3\u4e1a\u52a1\u8868\u4e0e\u8d44\u91d1\u6258\u7ba1\u6a21\u578b",
        left_title="\u4e1a\u52a1\u6570\u636e\u8868",
        left_content=[
            "sys_user \u2014 \u7528\u6237\u8868\uff08\u4f4d\u63a9\u7801\u89d2\u8272\uff09",
            "project \u2014 \u9879\u76ee\u8868\uff08\u72b6\u6001\u673a\uff09",
            "support_order \u2014 \u652f\u6301\u8ba2\u5355\u8868",
            "project_reward \u2014 \u652f\u6301\u6863\u4f4d\u4e0e\u56de\u62a5",
            "user_comment \u2014 \u8bc4\u8bba\u4e0e\u56de\u590d",
            "sys_notice \u2014 \u5e73\u53f0\u516c\u544a",
            "sys_user_message \u2014 \u7528\u6237\u6d88\u606f",
            "admin_operation_log \u2014 \u7ba1\u7406\u5458\u64cd\u4f5c\u65e5\u5fd7",
        ],
        right_title="\u8d44\u91d1\u6258\u7ba1\u8868",
        right_content=[
            "fund_account \u2014 \u9879\u76ee\u8d44\u91d1\u6258\u7ba1\u8d26\u6237",
            "  \u2022 total/frozen/available/paid/refunded",
            "fund_ledger \u2014 \u8d44\u91d1\u4ea4\u6613\u6d41\u6c34\u8868",
            "  \u2022 \u7c7b\u578b\uff1a\u5165\u8d26/\u9000\u6b3e/\u62e8\u4ed8",
            "fund_payout_phase \u2014 \u9636\u6bb5\u62e8\u4ed8\u8ba1\u5212",
            "  \u2022 \u6bd4\u4f8b\u62e8\u4ed8 + \u5ba1\u6838\u72b6\u6001",
            "project_payout \u2014 \u62e8\u4ed8\u8bb0\u5f55",
            "withdrawal_order \u2014 \u63d0\u73b0\u8ba2\u5355",
        ],
        page_num=5, total_pages=total_pages
    )

    # ===== Slide 6: Module Overview =====
    gen.add_table_slide(
        title="\u529f\u80fd\u6a21\u5757\u603b\u89c8",
        subtitle="\u4e03\u5927\u6838\u5fc3\u6a21\u5757 \u00d7 \u56db\u7c7b\u89d2\u8272 \u00d7 59\u9879\u529f\u80fd\u9700\u6c42",
        headers=["\u6a21\u5757", "\u6838\u5fc3\u804c\u8d23", "\u9002\u7528\u89d2\u8272"],
        rows=[
            ["F1 \u8ba4\u8bc1\u4e0e\u7528\u6237\u7ba1\u7406", "\u6ce8\u518c\u767b\u5f55\u3001JWT\u9274\u6743\u3001\u5b9e\u540d\u8ba4\u8bc1\u3001\u652f\u4ed8\u7ed1\u5b9a", "\u5168\u89d2\u8272"],
            ["F2 \u9879\u76ee\u6d4f\u89c8\u4e0e\u4ea4\u4e92", "\u5927\u5385\u5c55\u793a\u3001\u591a\u7ef4\u68c0\u7d22\u3001\u8bc4\u8bba\u70b9\u8d5e\u3001\u8d8b\u52bf\u56fe\u8868", "\u6e38\u5ba2+\u652f\u6301\u8005"],
            ["F3 \u9879\u76ee\u652f\u6301\u4e0e\u4ea4\u6613", "\u6863\u4f4d\u652f\u6301\u3001\u591a\u6e20\u9053\u652f\u4ed8\u3001\u8ba2\u5355\u72b6\u6001\u673a", "\u652f\u6301\u8005"],
            ["F4 \u9879\u76ee\u53d1\u8d77\u4e0e\u7ba1\u7406", "\u521b\u5efa\u9879\u76ee\u3001\u5ba1\u6838\u6d41\u7a0b\u3001\u52a8\u6001\u53d1\u5e03\u3001\u6570\u636e\u770b\u677f", "\u53d1\u8d77\u4eba"],
            ["F5 \u540e\u53f0\u76d1\u7ba1\u4e0e\u8fd0\u8425", "\u7528\u6237\u7ba1\u7406\u3001\u9879\u76ee\u5ba1\u6838\u3001\u6570\u636e\u62a5\u8868\u3001\u516c\u544a\u53d1\u5e03", "\u7ba1\u7406\u5458"],
            ["F6 \u6d88\u606f\u4e0e\u901a\u77e5", "\u9879\u76ee\u52a8\u6001\u63a8\u9001\u3001\u516c\u544a\u5e7f\u64ad\u3001WebSocket\u5b9e\u65f6\u901a\u4fe1", "\u5168\u89d2\u8272"],
            ["F7 \u8de8\u7aef\u9002\u914d\u4e0e\u4f53\u9a8c", "\u54cd\u5e94\u5f0f\u5e03\u5c40\u3001\u79fb\u52a8\u7aef\u7ba1\u7406\u3001\u5b9e\u65f6\u6d88\u606f", "\u5168\u5e73\u53f0"],
        ],
        page_num=6, total_pages=total_pages
    )

    # ===== Slide 7: Auth & Permission =====
    gen.add_two_column_slide(
        title="\u8ba4\u8bc1\u4e0e\u6743\u9650\u6a21\u5757",
        subtitle="Spring Security + JWT + \u4f4d\u63a9\u7801RBAC",
        left_title="\u8ba4\u8bc1\u673a\u5236",
        left_content=[
            "Bcrypt\u5355\u5411\u54c8\u5e0c\u52a0\u5bc6\u5b58\u50a8\u5bc6\u7801",
            "JWT\u65e0\u72b6\u6001\u4ee4\u724c\uff0c\u643a\u5e26\u7528\u6237ID\u4e0e\u89d2\u8272\u4f4d\u63a9\u7801",
            "Authorization: Bearer <token> \u7edf\u4e00\u9274\u6743",
            "Redis\u7f13\u5b58\u9a8c\u8bc1\u7801\uff0c5\u5206\u949fTTL\u9632\u5237",
            "\u767b\u5f55/\u9000\u51fa/\u5bc6\u7801\u4fee\u6539/\u627e\u56de\u5bc6\u7801\u5b8c\u6574\u94fe\u8def",
        ],
        right_title="\u6743\u9650\u6a21\u578b",
        right_content=[
            "\u4f4d\u63a9\u7801\u89d2\u8272\uff1aUSER=1, SPONSOR=2, ADMIN=4",
            "\u6309\u4f4d\u4e0e\u8fd0\u7b97\u5224\u65ad\u591a\u89d2\u8272\u53e0\u52a0",
            "\u540c\u4e00\u8d26\u6237\u53ef\u540c\u65f6\u62e5\u6709\u591a\u4e2a\u89d2\u8272",
            "SpringSecurity\u62e6\u622a\u94fe\u7cbe\u51c6\u63a7\u5236\u63a5\u53e3\u8bbf\u95ee",
            "\u7ba1\u7406\u5458\u53ef\u52a8\u6001\u5206\u914d/\u4fee\u6539\u7528\u6237\u89d2\u8272",
        ],
        page_num=7, total_pages=total_pages
    )

    # ===== Slide 8: Fund Custody Model (KEY) =====
    gen.add_content_slide(
        title="\u9879\u76ee\u652f\u6301\u4e0e\u8d44\u91d1\u6258\u7ba1\u6a21\u578b",
        subtitle="\u672c\u7cfb\u7edf\u6838\u5fc3\u4eae\u70b9\uff1a\u6258\u7ba1\u8d26\u6237 + \u8d44\u91d1\u6d41\u6c34\u8d26\u672c\u53cc\u5c42\u5efa\u6a21",
        bullets=[
            "\u6258\u7ba1\u8d26\u6237\u6a21\u578b\uff1afund_account \u8bb0\u5f55\u9879\u76ee\u8d44\u91d1\u72b6\u6001\uff08\u603b\u989d/\u51bb\u7ed3/\u53ef\u7528/\u5df2\u62e8\u4ed8/\u5df2\u9000\u6b3e\uff09",
            "\u8d44\u91d1\u6d41\u6c34\u8d26\u672c\uff1afund_ledger \u8bb0\u5f55\u6bcf\u7b14\u4ea4\u6613\uff08\u5165\u8d26/\u9000\u6b3e/\u62e8\u4ed8\uff09\uff0c\u53ef\u8ffd\u6eaf\u53ef\u5ba1\u8ba1",
            "\u9636\u6bb5\u62e8\u4ed8\u673a\u5236\uff1afund_payout_phase \u6309\u6bd4\u4f8b\u5206\u9636\u6bb5\u62e8\u4ed8\uff0c\u7ba1\u7406\u5458\u5ba1\u6838\u540e\u89e3\u51bb",
            "\u4e8b\u52a1\u4fdd\u969c\uff1a@Transactional \u786e\u4fdd\u8ba2\u5355\u521b\u5efa\u3001\u91d1\u989d\u8ba1\u7b97\u3001\u8d26\u6237\u66f4\u65b0\u7684\u539f\u5b50\u6027",
            "\u9632\u91cd\u590d\u63d0\u4ea4\uff1aRedis\u5206\u5e03\u5f0f\u9501 order:lock:{userId}:{projectId}",
            "\u652f\u4ed8\u56de\u8c03\u9a8c\u7b7e\uff1aRSA\u7b7e\u540d\u9a8c\u8bc1\uff0c\u6821\u9a8c\u91d1\u989d\u3001\u8ba2\u5355\u53f7\u9632\u7be1\u6539",
        ],
        page_num=8, total_pages=total_pages
    )

    # ===== Slide 9: State Machine =====
    gen.add_two_column_slide(
        title="\u72b6\u6001\u673a\u8bbe\u8ba1",
        subtitle="\u9879\u76ee\u72b6\u6001\u673a\u4e0e\u8ba2\u5355\u72b6\u6001\u673a\u7684\u5f3a\u4e00\u81f4\u6027\u4fdd\u969c",
        left_title="\u9879\u76ee\u72b6\u6001\u673a",
        left_content=[
            "\u5f85\u5ba1\u6838 (PENDING_REVIEW)",
            "  \u2193 \u5ba1\u6838\u901a\u8fc7",
            "\u7b79\u6b3e\u4e2d (FUNDING)",
            "  \u2193 \u7b79\u6b3e\u6210\u529f / \u5230\u671f",
            "\u5df2\u5b8c\u6210 (SUCCESS) / \u5df2\u53d6\u6d88",
            "\u5f02\u5e38\u5206\u652f\uff1a\u5df2\u9a73\u56de / \u5df2\u4e0b\u67b6",
            "\u540e\u7aef\u5f3a\u5236\u6821\u9a8c\u72b6\u6001\u6d41\u8f6c\u5408\u6cd5\u6027",
        ],
        right_title="\u8ba2\u5355\u72b6\u6001\u673a",
        right_content=[
            "\u5f85\u652f\u4ed8 (PENDING_PAYMENT)",
            "  \u2193 \u652f\u4ed8\u6210\u529f",
            "\u5df2\u652f\u4ed8 (PAID)",
            "  \u2193 \u9879\u76ee\u5931\u8d25",
            "\u5df2\u9000\u6b3e (REFUNDED)",
            "\u5f02\u5e38\u5206\u652f\uff1a\u5df2\u53d6\u6d88",
            "\u72b6\u6001\u53d8\u66f4\u540c\u6b65\u89e6\u53d1\u8d44\u91d1\u6d41\u6c34\u8bb0\u5f55",
        ],
        page_num=9, total_pages=total_pages
    )

    # ===== Slide 10: Admin Dashboard =====
    gen.add_content_slide(
        title="\u540e\u53f0\u76d1\u7ba1\u4e0e\u6570\u636e\u53ef\u89c6\u5316",
        subtitle="\u7ba1\u7406\u5458\u5168\u6743\u9650\u8fd0\u8425\u4e2d\u53f0",
        bullets=[
            "\u7528\u6237\u76d1\u7ba1\uff1a\u5168\u91cf\u7528\u6237\u5217\u8868\u3001\u89d2\u8272\u5206\u914d\u3001\u8d26\u53f7\u7981\u7528\u3001\u64cd\u4f5c\u65e5\u5fd7\u7559\u75d5",
            "\u9879\u76ee\u5ba1\u6838\u4e0e\u98ce\u63a7\uff1a\u5f85\u5ba1\u6838\u5217\u8868\u3001\u901a\u8fc7/\u9a73\u56de\u64cd\u4f5c\u3001\u8fdd\u89c4\u9879\u76ee\u5f3a\u5236\u4e0b\u67b6",
            "\u6570\u636e\u770b\u677f\uff1aECharts\u6e32\u67d3\u591a\u7ef4\u56fe\u8868\uff08\u67f1\u72b6\u56fe/\u6298\u7ebf\u56fe/\u997c\u56fe\uff09",
            "\u62a5\u8868\u5bfc\u51fa\uff1aEasyExcel\u4f4e\u5185\u5b58\u6d41\u5f0f\u5199\u5165.xlsx\u6587\u4ef6",
            "\u516c\u544a\u7cfb\u7edf\uff1a\u53d1\u5e03\u5e73\u53f0\u516c\u544a\u3001\u5b9a\u5411/\u5168\u91cf\u63a8\u9001",
            "\u63d0\u73b0\u5ba1\u6838\uff1a\u53d1\u8d77\u4eba\u63d0\u73b0\u7533\u8bf7\u3001\u7ba1\u7406\u5458\u5ba1\u6838\u62e8\u4ed8",
        ],
        page_num=10, total_pages=total_pages
    )

    # ===== Slide 11: WebSocket =====
    gen.add_content_slide(
        title="\u5b9e\u65f6\u6d88\u606f\u63a8\u9001",
        subtitle="WebSocket\u5168\u53cc\u5de5\u901a\u4fe1 + \u4e8b\u4ef6\u9a71\u52a8\u63a8\u9001",
        bullets=[
            "spring-boot-starter-websocket\u96c6\u6210\uff0c\u81ea\u5b9a\u4e49ServerEndpointConfig",
            "\u63e1\u624b\u9636\u6bb5\u901a\u8fc7JWT\u6821\u9a8c\uff0c\u4ec5\u5408\u6cd5Token\u53ef\u5efa\u7acbws\u8fde\u63a5",
            "ConcurrentHashMap<String, Session>\u7ef4\u62a4\u5728\u7ebf\u7528\u6237\u6c60",
            "\u4e8b\u4ef6\u89e6\u53d1\u673a\u5236\uff1a\u9879\u76ee\u72b6\u6001\u53d8\u66f4/\u516c\u544a\u53d1\u5e03\u65f6\u4e3b\u52a8push\u7ed3\u6784\u5316JSON",
            "\u6d88\u606f\u7c7b\u578b\uff1a\u9879\u76ee\u52a8\u6001\u3001\u8ba2\u5355\u72b6\u6001\u3001\u5e73\u53f0\u516c\u544a\u3001\u8bc4\u8bba\u56de\u590d",
            "\u79fb\u52a8\u7aef\u540c\u6837\u652f\u6301WebSocket\u5b9e\u65f6\u63a5\u6536\u63a8\u9001\u63d0\u9192",
        ],
        page_num=11, total_pages=total_pages
    )

    # ===== Slide 12: Key Technical Highlights =====
    gen.add_highlight_slide(
        title="\u5173\u952e\u6280\u672f\u4eae\u70b9",
        subtitle="\u4f17\u7b79\u573a\u666f\u7279\u6709\u7684\u5de5\u7a0b\u5316\u89e3\u51b3\u65b9\u6848",
        highlights=[
            ("\u8d44\u91d1\u6258\u7ba1\u6a21\u578b",
             "\u6258\u7ba1\u8d26\u6237+\u6d41\u6c34\u8d26\u672c\u53cc\u5c42\u5efa\u6a21\uff0c\u8d44\u91d1\u9694\u79bb\u3001\u53ef\u8ffd\u6eaf\u3001\u53ef\u5ba1\u8ba1\uff0c\u9636\u6bb5\u62e8\u4ed8\u673a\u5236\u4fdd\u969c\u8d44\u91d1\u5b89\u5168"),
            ("\u4f4d\u63a9\u7801RBAC",
             "\u7528\u6237\u89d2\u8272\u4ee5\u6574\u578b\u4f4d\u63a9\u7801\u4fdd\u5b58\uff0c\u6309\u4f4d\u4e0e\u8fd0\u7b97\u5224\u65ad\u6743\u9650\uff0c\u652f\u6301\u591a\u89d2\u8272\u53e0\u52a0\uff0c\u7075\u6d3b\u53ef\u6269\u5c55"),
            ("\u5206\u5e03\u5f0f\u9501\u9632\u5e76\u53d1",
             "Redis SET NX EX\u539f\u8bed\u5b9e\u73b0\u5206\u5e03\u5f0f\u9501\uff0c\u9632\u6b62\u8ba2\u5355\u8d85\u5356\u3001\u91cd\u590d\u63d0\u4ea4\uff0c\u4fdd\u969c\u4ea4\u6613\u539f\u5b50\u6027"),
            ("\u5207\u9762\u5f0f\u9650\u6d41",
             "@RateLimit\u81ea\u5b9a\u4e49\u6ce8\u89e3+AOP\u5207\u9762\uff0cRedis INCR\u56fa\u5b9a\u7a97\u53e3\u8ba1\u6570\uff0c\u63a5\u53e3\u7ea7\u7cbe\u7ec6\u5316\u9650\u6d41\u63a7\u5236"),
            ("\u72b6\u6001\u673a\u6cbb\u7406",
             "\u9879\u76ee/\u8ba2\u5355\u72b6\u6001\u679a\u4e3e\u7c7b\u5f3a\u5236\u7ea6\u675f\uff0c\u540e\u7aef\u6821\u9a8c\u6d41\u8f6c\u5408\u6cd5\u6027\uff0c\u4e8b\u52a1\u4fdd\u969c\u5f3a\u4e00\u81f4\u6027"),
            ("\u5168\u94fe\u8def\u5b89\u5168",
             "Bcrypt\u54c8\u5e0c+JWT\u65e0\u72b6\u6001+RSA\u56de\u8c03\u9a8c\u7b7e+XSS\u8fc7\u6ee4+SQL\u9884\u7f16\u8bd1\uff0c\u591a\u5c42\u6b21\u5b89\u5168\u9632\u62a4"),
        ],
        page_num=12, total_pages=total_pages
    )

    # ===== Slide 13: Testing =====
    gen.add_two_column_slide(
        title="\u7cfb\u7edf\u6d4b\u8bd5",
        subtitle="\u529f\u80fd\u6d4b\u8bd5 + \u63a5\u53e3\u538b\u6d4b + \u517c\u5bb9\u6027 + \u5b89\u5168\u6d4b\u8bd5",
        left_title="\u6d4b\u8bd5\u7c7b\u578b\u4e0e\u8303\u56f4",
        left_content=[
            "\u529f\u80fd\u6d4b\u8bd5\uff1a59\u9879\u529f\u80fd\u9700\u6c42\u5168\u8986\u76d6\u9a8c\u8bc1",
            "\u63a5\u53e3\u6d4b\u8bd5\uff1aKnife4j\u6587\u6863+Postman\u5192\u70df\u6d4b\u8bd5",
            "\u6027\u80fd\u6d4b\u8bd5\uff1a\u6838\u5fc3\u63a5\u53e3\u538b\u6d4b\uff0c\u9a8c\u8bc1\u5e76\u53d1\u573a\u666f\u7a33\u5b9a\u6027",
            "\u517c\u5bb9\u6027\u6d4b\u8bd5\uff1aPC/\u5e73\u677f/\u624b\u673a\u591a\u7aef\u9002\u914d\u9a8c\u8bc1",
            "\u5b89\u5168\u6d4b\u8bd5\uff1a\u8d8a\u6743\u8bbf\u95ee\u3001SQL\u6ce8\u5165\u3001XSS\u653b\u51fb\u6d4b\u8bd5",
        ],
        right_title="\u6d4b\u8bd5\u7ed3\u679c\u6458\u8981",
        right_content=[
            "\u529f\u80fd\u6d4b\u8bd5\u901a\u8fc7\u7387\uff1a100%",
            "\u6838\u5fc3\u63a5\u53e3\u54cd\u5e94\u65f6\u95f4 < 200ms",
            "\u5e76\u53d1\u573a\u666f\u4e0b\u65e0\u8d44\u91d1\u5f02\u5e38\u3001\u65e0\u8d85\u5356",
            "\u591a\u7aef\u5e03\u5c40\u65e0\u9519\u4e71\u3001\u529f\u80fd\u6b63\u5e38",
            "\u8d8a\u6743\u8bbf\u95ee\u88ab\u62e6\u622a\uff0c\u6ce8\u5165/\u6ce8\u5165\u653b\u51fb\u88ab\u9632\u5fa1",
        ],
        page_num=13, total_pages=total_pages
    )

    # ===== Slide 14: System Demo =====
    gen.add_content_slide(
        title="\u7cfb\u7edf\u6f14\u793a\u4e0e\u754c\u9762\u5c55\u793a",
        subtitle="\u524d\u7aef\u754c\u9762\u4e0e\u6838\u5fc3\u4ea4\u4e92\u6d41\u7a0b",
        bullets=[
            "\u9996\u9875\u5927\u5385\uff1a\u70ed\u95e8\u9879\u76ee\u6392\u884c\u699c + \u591a\u7ef4\u7b5b\u9009 + \u4e2a\u6027\u5316\u63a8\u8350",
            "\u9879\u76ee\u8be6\u60c5\uff1a\u89c6\u9891\u64ad\u653e + \u7b79\u6b3e\u8d8b\u52bf\u56fe + \u6863\u4f4d\u652f\u6301 + \u8bc4\u8bba\u4e92\u52a8",
            "\u53d1\u8d77\u4eba\u540e\u53f0\uff1a\u9879\u76ee\u7ba1\u7406 + \u6570\u636e\u770b\u677f + \u652f\u6301\u8005\u5bfc\u51fa",
            "\u7ba1\u7406\u5458\u4e2d\u53f0\uff1a\u5ba1\u6838\u5217\u8868 + \u6570\u636e\u53ef\u89c6\u5316\u56fe\u8868 + \u62a5\u8868\u5bfc\u51fa",
            "\u4e2a\u4eba\u4e2d\u5fc3\uff1a\u8d44\u6599\u7ef4\u62a4 + \u5b9e\u540d\u8ba4\u8bc1 + \u652f\u4ed8\u7ed1\u5b9a + \u6211\u7684\u652f\u6301",
            "\u79fb\u52a8\u7aef\u9002\u914d\uff1a\u54cd\u5e94\u5f0f\u5e03\u5c40\u3001\u5168\u529f\u80fd\u53ef\u7528",
        ],
        page_num=14, total_pages=total_pages
    )

    # ===== Slide 15: Conclusion =====
    gen.add_summary_slide(
        title="\u603b\u7ed3\u4e0e\u5c55\u671b",
        points=[
            "\u5b8c\u6210\u4e86\u57fa\u4e8eSpring Boot + Vue 3\u7684\u5168\u6808\u4f17\u7b79\u5e73\u53f0FundTogether\u7684\u8bbe\u8ba1\u4e0e\u5b9e\u73b0",
            "\u5b9e\u73b0\u4e864\u7c7b\u89d2\u8272\u00d77\u5927\u6a21\u5757\u00d759\u9879\u529f\u80fd\u9700\u6c42\u7684\u5b8c\u6574\u4e1a\u52a1\u95ed\u73af",
            "\u63d0\u51fa\u5e76\u5b9e\u73b0\u4e86\u201c\u6258\u7ba1\u8d26\u6237+\u8d44\u91d1\u6d41\u6c34\u201d\u53cc\u5c42\u5efa\u6a21\u7684\u8d44\u91d1\u5b89\u5168\u65b9\u6848",
            "\u7cfb\u7edf\u7ecf\u8fc7\u529f\u80fd\u6d4b\u8bd5\u3001\u63a5\u53e3\u538b\u6d4b\u3001\u5b89\u5168\u6d4b\u8bd5\uff0c\u5404\u6a21\u5757\u8fd0\u884c\u7a33\u5b9a",
        ],
        conclusion="\u672a\u6765\u5c55\u671b\uff1a\u5fae\u670d\u52a1\u5316\u62c6\u5206 \u00b7 \u5f39\u6027\u6269\u5c55\u4e0e\u8d1f\u8f7d\u5747\u8861 \u00b7 \u667a\u80fd\u63a8\u8350\u7b97\u6cd5\u4f18\u5316 \u00b7 \u79fb\u52a8\u7aef\u5c0f\u7a0b\u5e8f\u62d3\u5c55",
        page_num=15, total_pages=total_pages
    )

    # ===== Thank You =====
    gen.add_thank_you_slide()

    output_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "..",
        "\u6bd5\u4e1a\u7b54\u8fa9_PPT_FundTogether.pptx"
    )
    output_path = os.path.abspath(output_path)
    gen.save(output_path)
    print(f"PPT generated successfully: {output_path}")
    return output_path


if __name__ == "__main__":
    build_presentation()
